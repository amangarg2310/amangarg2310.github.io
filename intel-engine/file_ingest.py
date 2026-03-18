"""
File ingestion — extract text from PDF, DOCX, and PPTX documents.
"""

import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_pdf_text(file_path: str) -> tuple[str, int]:
    """Extract text from a PDF file. Returns (text, page_count)."""
    import pdfplumber

    text_parts = []
    with pdfplumber.open(file_path) as pdf:
        page_count = len(pdf.pages)
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)

    text = "\n\n".join(text_parts)
    if not text.strip():
        raise ValueError("PDF contains no extractable text (may be image-based)")

    logger.info(f"Extracted {len(text)} chars from {page_count}-page PDF")
    return text, page_count


def extract_docx_text(file_path: str) -> tuple[str, int]:
    """Extract text from a DOCX file. Returns (text, paragraph_count)."""
    from docx import Document

    doc = Document(file_path)
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)

    # Also extract from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)

    text = "\n\n".join(paragraphs)
    if not text.strip():
        raise ValueError("DOCX contains no extractable text")

    logger.info(f"Extracted {len(text)} chars from DOCX ({len(paragraphs)} paragraphs)")
    return text, len(paragraphs)


def extract_pptx_text(file_path: str) -> tuple[str, int]:
    """Extract text from a PPTX file. Returns (text, slide_count)."""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides_text = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        slide_parts.append(row_text)

        if len(slide_parts) > 1:  # More than just the header
            slides_text.append("\n".join(slide_parts))

    slide_count = len(prs.slides)
    text = "\n\n".join(slides_text)
    if not text.strip():
        raise ValueError("PPTX contains no extractable text")

    logger.info(f"Extracted {len(text)} chars from {slide_count}-slide PPTX")
    return text, slide_count


def ingest_file(file_path: str, original_filename: str) -> tuple[str, int]:
    """
    Extract text from a document file based on its extension.

    Returns (text_content, page_or_element_count).
    """
    ext = Path(original_filename).suffix.lower()

    if ext == ".pdf":
        return extract_pdf_text(file_path)
    elif ext == ".docx":
        return extract_docx_text(file_path)
    elif ext == ".pptx":
        return extract_pptx_text(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
